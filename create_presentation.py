#!/usr/bin/env python3
"""
Generate: Overview Operasional dan Proses Produksi Pabrik Gula
Professional PowerPoint Presentation - 21 Slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# ============================================================
# COLOR PALETTE
# ============================================================
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
SUGARCANE_GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INDUSTRIAL_GRAY = RGBColor(0x42, 0x42, 0x42)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x9E, 0x9E, 0x9E)
WARM_BROWN = RGBColor(0x6D, 0x4C, 0x41)
GOLD_ACCENT = RGBColor(0xF9, 0xA8, 0x25)
DARK_BG = RGBColor(0x1A, 0x23, 0x2F)
CARD_BG = RGBColor(0xE8, 0xF5, 0xE9)

# ============================================================
# HELPER FUNCTIONS
# ============================================================


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_text(slide, left, top, width, height, text,
                        font_size=12, bold=False, color=INDUSTRIAL_GRAY,
                        fill_color=None, align=PP_ALIGN.LEFT,
                        font_name='Calibri', shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add a shape with text to a slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return shape



def add_textbox(slide, left, top, width, height, text,
                font_size=12, bold=False, color=INDUSTRIAL_GRAY,
                align=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multi_text(slide, left, top, width, height, lines,
                   font_size=11, color=INDUSTRIAL_GRAY, font_name='Calibri',
                   line_spacing=1.2, bold_first=False, align=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox



def add_footer(slide, slide_num, total=21):
    """Add footer with slide number and branding."""
    footer_text = f"{slide_num} / {total}  •  Overview Operasional & Proses Produksi Pabrik Gula  |  Internal Use"
    add_textbox(slide, Cm(1), Cm(17.5), Cm(30), Cm(1),
                footer_text, font_size=8, color=MID_GRAY, align=PP_ALIGN.LEFT)


def add_section_header(slide, title, subtitle=""):
    """Add consistent section header."""
    add_textbox(slide, Cm(2), Cm(1.5), Cm(28), Cm(2),
                title, font_size=28, bold=True, color=DARK_GREEN,
                font_name='Calibri', align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Cm(2), Cm(3.5), Cm(28), Cm(1.2),
                    subtitle, font_size=13, color=SUGARCANE_GREEN,
                    font_name='Calibri', align=PP_ALIGN.LEFT)


def add_key_takeaway(slide, text, top=Cm(15.5)):
    """Add a key takeaway box at bottom of slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), top, Cm(29), Cm(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = SUGARCANE_GREEN
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = f"Key Takeaway: {text}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.font.name = 'Calibri'


def add_process_box(slide, left, top, width, height, number, title, desc, accent_color=SUGARCANE_GREEN):
    """Add a numbered process box."""
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Cm(1.2), Cm(1.2))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent_color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Text box
    add_multi_text(slide, left + Cm(1.5), top, width - Cm(1.5), height,
                   [title, desc], font_size=10, bold_first=True, color=INDUSTRIAL_GRAY)



# ============================================================
# CREATE PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = Cm(33.867)  # 16:9 widescreen
prs.slide_height = Cm(19.05)

# Use blank layout
blank_layout = prs.slide_layouts[6]

# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_GREEN)

# Title
add_textbox(slide, Cm(3), Cm(2), Cm(28), Cm(1.5),
            "PRESENTASI INTERNAL  •  2026", font_size=11,
            color=GOLD_ACCENT, align=PP_ALIGN.LEFT)

add_textbox(slide, Cm(3), Cm(4), Cm(28), Cm(3),
            "Overview Operasional &\nProses Produksi Pabrik Gula",
            font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_textbox(slide, Cm(3), Cm(8.5), Cm(28), Cm(1.5),
            "Dari Penerimaan Tebu hingga Produk Akhir — Alur, GMP, QC & Sustainability",
            font_size=14, color=RGBColor(0xA5, 0xD6, 0xA7), align=PP_ALIGN.LEFT)

# Decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(11), Cm(8), Cm(0.15))
line.fill.solid()
line.fill.fore_color.rgb = GOLD_ACCENT
line.line.fill.background()

add_textbox(slide, Cm(3), Cm(12), Cm(28), Cm(2),
            "Disusun untuk pemahaman operasional pabrik gula\nBerbasis referensi: Hugot (1960), GMP Manual, Laporan KP & Jurnal Industri",
            font_size=10, color=RGBColor(0xC8, 0xE6, 0xC9), align=PP_ALIGN.LEFT)



# ============================================================
# SLIDE 2: DAFTAR ISI
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Daftar Isi", "Struktur Presentasi — 21 Bagian")
add_footer(slide, 2)

toc_items = [
    ("01", "Pendahuluan Industri Gula"),
    ("02", "Overview Alur Produksi"),
    ("03", "Penerimaan & Penanganan Tebu"),
    ("04", "Cane Preparation & Milling"),
    ("05", "Pemrosesan Nira (Clarification)"),
    ("06", "Evaporation Process"),
    ("07", "Crystallization Process"),
    ("08", "Centrifugation & Sugar Drying"),
    ("09", "Packaging & Storage"),
    ("10", "Utility Systems"),
    ("11", "Pemanfaatan Produk Samping"),
    ("12", "Good Manufacturing Practice"),
    ("13", "Quality Control"),
    ("14", "Pengelolaan Limbah"),
    ("15", "Sustainability & Efisiensi"),
    ("16", "Keselamatan Kerja (K3)"),
    ("17", "Tantangan Operasional"),
    ("18", "Kesimpulan"),
]

col1_items = toc_items[:9]
col2_items = toc_items[9:]

for i, (num, title) in enumerate(col1_items):
    y = Cm(5) + Cm(i * 1.3)
    add_textbox(slide, Cm(3), y, Cm(2), Cm(1),
                num, font_size=11, bold=True, color=SUGARCANE_GREEN)
    add_textbox(slide, Cm(5.5), y, Cm(12), Cm(1),
                title, font_size=11, color=INDUSTRIAL_GRAY)

for i, (num, title) in enumerate(col2_items):
    y = Cm(5) + Cm(i * 1.3)
    add_textbox(slide, Cm(18), y, Cm(2), Cm(1),
                num, font_size=11, bold=True, color=SUGARCANE_GREEN)
    add_textbox(slide, Cm(20.5), y, Cm(12), Cm(1),
                title, font_size=11, color=INDUSTRIAL_GRAY)



# ============================================================
# SLIDE 3: INTRODUCTION TO SUGAR INDUSTRY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Pendahuluan: Industri Gula Indonesia",
                   "Posisi strategis & konteks operasional")
add_footer(slide, 3)

# Key stats boxes
stats = [
    ("~2,2 jt ton", "Produksi nasional/tahun"),
    ("~5-6 jt ton", "Konsumsi nasional/tahun"),
    ("~50%", "Ketergantungan impor"),
    ("~60 PG", "Pabrik gula aktif"),
]
for i, (val, desc) in enumerate(stats):
    x = Cm(2.5) + Cm(i * 7.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(5), Cm(7), Cm(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = SUGARCANE_GREEN
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = INDUSTRIAL_GRAY
    p2.alignment = PP_ALIGN.CENTER

# Context points
context_lines = [
    "• Musim giling: April – November (6–7 bulan operasi)",
    "• Pemain utama: PTPN (BUMN), swasta nasional, pabrik gula rakyat",
    "• Komoditas strategis: pangan, devisa negara, lapangan kerja",
    "• Lead time tebu → gula: 24–48 jam (proses kontinu non-stop)",
    "• Kapasitas pabrik: 3.000 – 6.000 TCD (Ton Cane per Day)",
]
add_multi_text(slide, Cm(2.5), Cm(9), Cm(29), Cm(5),
               context_lines, font_size=11, color=INDUSTRIAL_GRAY)

add_key_takeaway(slide, "Pabrik gula adalah industri terintegrasi hulu-hilir dengan gap besar antara produksi dan konsumsi nasional.")



# ============================================================
# SLIDE 4: OVERVIEW ALUR PRODUKSI
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Overview Alur Produksi Pabrik Gula",
                   "9 Tahapan Utama — Dari Tebu Masuk hingga Gula Terkemas")
add_footer(slide, 4)

stages = [
    ("1", "Penerimaan\nTebu", SUGARCANE_GREEN),
    ("2", "Cane\nHandling", SUGARCANE_GREEN),
    ("3", "Milling\n(Gilingan)", DARK_GREEN),
    ("4", "Pemurnian\nNira", DARK_GREEN),
    ("5", "Evaporasi", DARK_GREEN),
    ("6", "Kristalisasi", WARM_BROWN),
    ("7", "Centrifugasi", WARM_BROWN),
    ("8", "Pengeringan\n& Pendinginan", WARM_BROWN),
    ("9", "Pengemasan\n& Gudang", GOLD_ACCENT),
]

for i, (num, label, col) in enumerate(stages):
    x = Cm(1.5) + Cm(i * 3.5)
    # Box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Cm(5.5), Cm(3.2), Cm(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = col
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    # Arrow (except last)
    if i < len(stages) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + Cm(3.2), Cm(6.5), Cm(0.4), Cm(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

# Categorization labels
add_textbox(slide, Cm(1.5), Cm(9), Cm(7.5), Cm(1),
            "PERSIAPAN BAHAN BAKU", font_size=9, bold=True, color=SUGARCANE_GREEN)
add_textbox(slide, Cm(9), Cm(9), Cm(12), Cm(1),
            "PENGOLAHAN INTI", font_size=9, bold=True, color=DARK_GREEN)
add_textbox(slide, Cm(21), Cm(9), Cm(11), Cm(1),
            "FINISHING & OUTPUT", font_size=9, bold=True, color=WARM_BROWN)

# Summary
summary_lines = [
    "Apa yang terjadi di proses ini?",
    "Tebu segar ditimbang → dicacah → digiling untuk diambil niranya → nira dimurnikan →",
    "air diuapkan → gula dikristalkan → dipisahkan → dikeringkan → dikemas siap jual.",
    "",
    "Analogi sederhana: 100 ton tebu masuk → keluar 8–10 ton gula + 30 ton ampas (bagasse) + 4–5 ton tetes (molasses)"
]
add_multi_text(slide, Cm(2), Cm(10.5), Cm(29), Cm(5),
               summary_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

add_key_takeaway(slide, "Total proses berjalan kontinu 24/7 selama musim giling. Setiap jam keterlambatan = kehilangan 0,1% pol.")



# ============================================================
# SLIDE 5: PENERIMAAN & PENANGANAN TEBU
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Penerimaan & Penanganan Tebu",
                   "Stasiun Pertama — Gerbang kualitas bahan baku")
add_footer(slide, 5)

# Process flow
flow_items = [
    ("Truk Masuk", "Tebu dari kebun tiba di pabrik"),
    ("Timbangan", "Ditimbang bruto & tarra"),
    ("Sampling", "Cek pol, brix, kadar trash"),
    ("Cane Yard", "Antri di emplasement (maks 24 jam)"),
    ("Feeding Table", "Dipindah ke meja tebu menuju gilingan"),
]

for i, (title, desc) in enumerate(flow_items):
    x = Cm(2) + Cm(i * 6.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Cm(5), Cm(5.8), Cm(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else WHITE
    shape.line.color.rgb = SUGARCANE_GREEN
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = f"⟶ {title}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = INDUSTRIAL_GRAY
    p2.alignment = PP_ALIGN.CENTER

# Explanation
explain_lines = [
    "Apa yang terjadi di proses ini?",
    "",
    "• Tebu dari kebun diangkut truk, ditimbang untuk mengetahui berat bersih",
    "• Dilakukan sampling: diukur kadar gula (pol), padatan terlarut (brix), dan kadar kotoran (trash)",
    "• Hasil sampling menentukan kualitas tebu & pembayaran ke petani",
    "• Tebu menunggu di cane yard (emplasement) — TIDAK BOLEH >24 jam karena pol turun",
    "• Sistem FIFO (First In First Out) untuk menjaga kesegaran tebu",
    "",
    "Parameter kritis: Pol tebu ≥8%, Trash <5%, Waktu tunggu <24 jam"
]
add_multi_text(slide, Cm(2), Cm(8.5), Cm(29), Cm(7),
               explain_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

add_key_takeaway(slide, "Kualitas tebu di titik ini menentukan rendemen akhir. Tebu terlambat = gula hilang.")



# ============================================================
# SLIDE 6: CANE PREPARATION & MILLING
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Cane Preparation & Milling",
                   "Stasiun Gilingan — Jantung ekstraksi gula")
add_footer(slide, 6)

# Two columns: Preparation & Milling
# Left: Preparation
add_textbox(slide, Cm(2), Cm(5), Cm(14), Cm(1),
            "CANE PREPARATION", font_size=12, bold=True, color=DARK_GREEN)
prep_lines = [
    "1. Cane Leveller — meratakan tebu di conveyor",
    "2. Cane Cutter (I & II) — memotong tebu jadi potongan pendek",
    "3. Unigrator / Shredder — mencacah tebu jadi serat halus",
    "",
    "Tujuan: membuka sel-sel tebu agar nira mudah diekstraksi.",
    "Preparation Index (PI) target: >85%",
    "(PI = seberapa banyak sel tebu yang terbuka)"
]
add_multi_text(slide, Cm(2), Cm(6.2), Cm(14), Cm(6),
               prep_lines, font_size=10, color=INDUSTRIAL_GRAY)

# Right: Milling
add_textbox(slide, Cm(17.5), Cm(5), Cm(14), Cm(1),
            "MILLING (GILINGAN)", font_size=12, bold=True, color=DARK_GREEN)
mill_lines = [
    "• Terdiri dari 4–6 unit gilingan (tandem mill)",
    "• Setiap unit: 3 roll (top, feed, discharge)",
    "• Tebu dicacah diperas bertahap dari gilingan I → terakhir",
    "• Ditambahkan air imbibisi di gilingan akhir",
    "  (untuk melarutkan sisa gula dalam ampas)",
    "",
    "OUTPUT:",
    "  → Nira mentah (mixed juice): brix ~12–15%",
    "  → Ampas/Bagasse: kadar air ~50%, jadi bahan bakar boiler",
    "",
    "Target: Pol extraction >95%, Moisture bagasse <50%"
]
add_multi_text(slide, Cm(17.5), Cm(6.2), Cm(15), Cm(8),
               mill_lines, font_size=10, color=INDUSTRIAL_GRAY)

add_key_takeaway(slide, "Milling adalah jantung pabrik — ekstraksi optimal di sini menentukan recovery pol keseluruhan (target >95%).")



# ============================================================
# SLIDE 7: PEMROSESAN NIRA (CLARIFICATION)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Pemrosesan Nira (Clarification)",
                   "Stasiun Pemurnian — Membersihkan nira dari pengotor")
add_footer(slide, 7)

# Process steps
clarif_steps = [
    ("Timbangan Nira", "Nira mentah ditimbang untuk kontrol massa"),
    ("Pemanasan I", "Dipanaskan 70–75°C (juice heater)"),
    ("Defekasi", "Ditambah susu kapur Ca(OH)₂ → pH naik ~7–8"),
    ("Sulfitasi", "Ditambah gas SO₂ → pH turun ~6,8–7,2\n(membantu pemucatan warna)"),
    ("Pemanasan II", "Dipanaskan 100–105°C"),
    ("Flash Tank", "Menghilangkan gas terlarut"),
    ("Clarifier", "Pengendapan kotoran (2–3 jam)\nHasil: nira jernih + endapan (mud)"),
    ("Rotary Vacuum Filter", "Mud disaring → blotong (filter cake)\nNira saring kembali ke proses"),
]

for i, (title, desc) in enumerate(clarif_steps):
    row = i // 4
    col = i % 4
    x = Cm(2) + Cm(col * 7.8)
    y = Cm(5) + Cm(row * 5.5)
    # Box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Cm(7.2), Cm(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG if row == 0 else WHITE
    shape.line.color.rgb = SUGARCANE_GREEN
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    tf.margin_left = Pt(8)
    p = tf.paragraphs[0]
    p.text = f"{i+1}. {title}"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = INDUSTRIAL_GRAY
    p2.space_before = Pt(4)

add_key_takeaway(slide, "Pemurnian menghasilkan nira jernih (clear juice) + blotong sebagai by-product pupuk organik.", Cm(16.5))



# ============================================================
# SLIDE 8: EVAPORATION PROCESS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Evaporation Process",
                   "Stasiun Penguapan — Menaikkan konsentrasi nira")
add_footer(slide, 8)

# Main explanation
evap_lines = [
    "Apa yang terjadi di proses ini?",
    "",
    "Nira jernih (brix ~12–15%) diuapkan airnya secara bertahap hingga",
    "menjadi nira kental / thick juice (brix ~60–65%).",
    "",
    "Sistem: Multiple Effect Evaporator (biasanya 4–5 badan/effect)",
    "",
    "Cara kerja:",
    "• Effect 1: dipanaskan uap bekas turbin (exhaust steam)",
    "• Effect 2–5: uap dari effect sebelumnya dipakai memanaskan effect berikutnya",
    "• Tekanan & suhu makin rendah di setiap effect (vacuum)",
    "",
    "Keuntungan sistem multi-effect:",
    "→ Hemat energi: 1 kg uap bisa menguapkan 4–5 kg air",
    "→ Suhu rendah di effect akhir: mencegah kerusakan gula (karamelisasi)",
]
add_multi_text(slide, Cm(2), Cm(5), Cm(18), Cm(11),
               evap_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

# Side info box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Cm(21.5), Cm(5), Cm(10.5), Cm(8))
shape.fill.solid()
shape.fill.fore_color.rgb = CARD_BG
shape.line.color.rgb = SUGARCANE_GREEN
shape.line.width = Pt(1)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_top = Pt(12)
tf.margin_left = Pt(10)
p = tf.paragraphs[0]
p.text = "Parameter Kunci"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = DARK_GREEN
params = [
    "\nInput: Brix 12–15%",
    "Output: Brix 60–65%",
    "Jumlah effect: 4–5",
    "Tekanan: vacuum",
    "Suhu effect akhir: ~55°C",
    "\nKontrol:",
    "• Brix nira kental",
    "• Suhu tiap effect",
    "• Vacuum pressure",
    "• Scaling (kerak)",
]
for line in params:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(9)
    p2.font.color.rgb = INDUSTRIAL_GRAY

add_key_takeaway(slide, "Evaporasi menghilangkan ~80% air dari nira. Efisiensi energi di sini sangat krusial untuk biaya produksi.")



# ============================================================
# SLIDE 9: CRYSTALLIZATION PROCESS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Crystallization Process",
                   "Stasiun Masakan — Pembentukan kristal gula")
add_footer(slide, 9)

cryst_lines = [
    "Apa yang terjadi di proses ini?",
    "",
    "Nira kental (thick juice) dimasak dalam vacuum pan pada suhu rendah",
    "dan tekanan vakum hingga terbentuk kristal gula (massecuite).",
    "",
    "Tahapan Masakan (3-Boiling System):",
    "",
    "• Masakan A (1st boiling):",
    "  - Nira kental + bibit kristal → kristal besar (gula produk SHS)",
    "  - Hasil: Massecuite A → centrifuge → Gula A + Molasses A",
    "",
    "• Masakan B (2nd boiling):",
    "  - Molasses A dimasak ulang → kristal sedang",
    "  - Hasil: Gula B (dilebur, kembali ke masakan A)",
    "",
    "• Masakan C/D (3rd boiling):",
    "  - Molasses B dimasak → kristal kecil (bibit/seed)",
    "  - Hasil: Gula C/D (bibit) + Final Molasses (tetes akhir)",
    "",
    "Parameter kritis:",
    "• Suhu masakan: 65–75°C (di bawah vakum)",
    "• Supersaturasi: dijaga ~1,1–1,2",
    "• Ukuran kristal: 0,8–1,1 mm untuk gula SHS",
]
add_multi_text(slide, Cm(2), Cm(5), Cm(20), Cm(12.5),
               cryst_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

# Side note
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Cm(23), Cm(5), Cm(9), Cm(6))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
shape.line.color.rgb = GOLD_ACCENT
shape.line.width = Pt(1)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_top = Pt(10)
tf.margin_left = Pt(10)
p = tf.paragraphs[0]
p.text = "Istilah Penting"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = WARM_BROWN
terms = [
    "\nMassecuite = campuran kristal + larutan induk",
    "Vacuum Pan = bejana masak bertekanan vakum",
    "Seed/Bibit = kristal kecil pemicu pertumbuhan",
    "Molasses = tetes, sisa cair setelah kristalisasi",
    "SHS = Superior High Sugar (pol ≥99,3%)",
]
for t in terms:
    p2 = tf.add_paragraph()
    p2.text = t
    p2.font.size = Pt(9)
    p2.font.color.rgb = INDUSTRIAL_GRAY

add_key_takeaway(slide, "Kristalisasi menentukan ukuran & kemurnian gula akhir. Kontrol supersaturasi sangat krusial.")



# ============================================================
# SLIDE 10: CENTRIFUGATION & SUGAR DRYING
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Centrifugation & Sugar Drying",
                   "Stasiun Putar & Pengeringan — Pemisahan dan finishing kristal")
add_footer(slide, 10)

# Left: Centrifugation
add_textbox(slide, Cm(2), Cm(5), Cm(14), Cm(1),
            "CENTRIFUGATION (PEMUTARAN)", font_size=12, bold=True, color=DARK_GREEN)
cent_lines = [
    "• Massecuite dari vacuum pan dimasukkan ke centrifuge",
    "• Mesin berputar 1000–1500 rpm",
    "• Kristal gula tertahan di saringan (screen)",
    "• Molasses (tetes) terlempar keluar melalui lubang screen",
    "• Dicuci dengan air/steam untuk membersihkan kristal",
    "",
    "Output: Kristal gula basah + molasses terpisah",
    "",
    "Tipe: Batch centrifuge (gula A) & Continuous (gula B/C)"
]
add_multi_text(slide, Cm(2), Cm(6.2), Cm(14), Cm(7),
               cent_lines, font_size=10, color=INDUSTRIAL_GRAY)

# Right: Drying
add_textbox(slide, Cm(17.5), Cm(5), Cm(14), Cm(1),
            "SUGAR DRYING & COOLING", font_size=12, bold=True, color=DARK_GREEN)
dry_lines = [
    "• Gula basah (moisture ~1%) masuk rotary dryer",
    "• Dikeringkan dengan udara panas (suhu ~80°C)",
    "• Target kadar air akhir: < 0,05% (gula SHS)",
    "• Setelah kering → masuk sugar cooler",
    "• Didinginkan ke suhu kamar (~30–35°C)",
    "• Tujuan: mencegah caking (penggumpalan) saat penyimpanan",
    "",
    "Quality check setelah drying:",
    "• Kadar air (moisture) < 0,05%",
    "• Warna (ICUMSA): 150–300 untuk GKP",
    "• Ukuran kristal: 0,8–1,1 mm",
    "• Pol ≥ 99,3%"
]
add_multi_text(slide, Cm(17.5), Cm(6.2), Cm(15), Cm(8),
               dry_lines, font_size=10, color=INDUSTRIAL_GRAY)

add_key_takeaway(slide, "Centrifuge memisahkan kristal dari molasses. Pengeringan memastikan gula tidak menggumpal saat disimpan.")



# ============================================================
# SLIDE 11: PACKAGING & STORAGE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Packaging & Storage",
                   "Stasiun Pengemasan — Dari bulk ke produk siap distribusi")
add_footer(slide, 11)

pack_lines = [
    "Apa yang terjadi di proses ini?",
    "",
    "Alur Pengemasan:",
    "1. Gula kering dari cooler → vibrating screen (ayakan) → pisahkan ukuran",
    "2. Gula sesuai spesifikasi → hopper penampung",
    "3. Penimbangan otomatis (50 kg/karung atau sesuai order)",
    "4. Pengemasan dalam karung (PP woven bag + inner PE liner)",
    "5. Dijahit/sealed → conveyor → gudang",
    "",
    "Penyimpanan (Gudang Gula):",
    "• Gudang harus kering, ventilasi baik, suhu <30°C",
    "• Kelembaban relatif (RH) < 65% — mencegah caking",
    "• Sistem FIFO (First In First Out)",
    "• Pallet kayu/plastik — tidak langsung kontak lantai",
    "• Penumpukan maksimum: 30–40 karung",
    "",
    "GMP di area packaging:",
    "• Area bersih, bebas hama, tertutup dari cuaca",
    "• Operator menggunakan APD (masker, sarung tangan, hairnet)",
    "• Label lengkap: tanggal produksi, batch, expired date, berat bersih",
]
add_multi_text(slide, Cm(2), Cm(5), Cm(29), Cm(12),
               pack_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

add_key_takeaway(slide, "Pengemasan & penyimpanan yang baik menjaga kualitas gula sampai ke tangan konsumen.")



# ============================================================
# SLIDE 12: UTILITY SYSTEMS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Utility Systems",
                   "Sistem Pendukung — Energi, air, dan infrastruktur operasional")
add_footer(slide, 12)

# Grid of utility systems
utilities = [
    ("Boiler / Ketel Uap", [
        "• Bahan bakar: bagasse (ampas tebu)",
        "• Menghasilkan uap bertekanan tinggi",
        "• Kapasitas: 20–60 ton uap/jam",
        "• Tipe: water tube boiler",
        "• Pabrik gula hampir self-sufficient energi"
    ]),
    ("Turbin & Power House", [
        "• Uap boiler → turbin → listrik + exhaust steam",
        "• Listrik: 3–8 MW (kebutuhan pabrik)",
        "• Exhaust steam: dipakai proses (evaporator, pan)",
        "• Co-generation: surplus dijual ke PLN",
        "• Back Pressure Turbine Generator (BPTG)"
    ]),
    ("Water Treatment Plant", [
        "• Sumber: sungai / sumur dalam",
        "• Treatment: koagulasi, filtrasi, softening",
        "• Kebutuhan: 500–1500 m³/hari",
        "• Untuk: boiler feed water, proses, sanitasi",
        "• Kualitas: hardness <5 ppm untuk boiler"
    ]),
    ("Instrument & Workshop", [
        "• Workshop mekanik & listrik",
        "• Spare parts management",
        "• Preventive maintenance scheduling",
        "• Laboratorium instrumen & kalibrasi",
        "• Off-season: overhaul seluruh mesin"
    ]),
]

for i, (title, items) in enumerate(utilities):
    row = i // 2
    col = i % 2
    x = Cm(2) + Cm(col * 15.8)
    y = Cm(5) + Cm(row * 5.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Cm(15), Cm(5.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG if (i % 2 == 0) else WHITE
    shape.line.color.rgb = SUGARCANE_GREEN
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(10)
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = item
        p2.font.size = Pt(9)
        p2.font.color.rgb = INDUSTRIAL_GRAY
        p2.space_before = Pt(2)

add_key_takeaway(slide, "Pabrik gula bisa mandiri energi dengan membakar bagasse di boiler — konsep self-sufficient & circular economy.")



# ============================================================
# SLIDE 13: PEMANFAATAN PRODUK SAMPING
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Pemanfaatan Produk Samping (By-Products)",
                   "Circular Economy — Tidak ada yang terbuang")
add_footer(slide, 13)

byproducts = [
    ("Bagasse (Ampas Tebu)", "~30 ton / 100 ton tebu",
     "• Bahan bakar utama boiler\n• Co-generation listrik\n• Potensi: pulp & kertas, pellet biofuel\n• Revenue: 5–10%",
     SUGARCANE_GREEN),
    ("Molasses (Tetes)", "~4–5 ton / 100 ton tebu",
     "• Bahan baku alkohol/etanol\n• Pakan ternak\n• MSG, ragi, asam sitrat\n• Revenue: 10–15%",
     WARM_BROWN),
    ("Blotong (Filter Cake)", "~3–4 ton / 100 ton tebu",
     "• Pupuk organik (N, P, K tinggi)\n• Soil conditioner kebun tebu\n• Kompos\n• Revenue: 2–3%",
     DARK_GREEN),
    ("Listrik Surplus (Co-gen)", "Kapasitas 3–8 MW",
     "• Dijual ke PLN (excess power)\n• Bisa mengurangi biaya energi\n• Green energy certificate\n• Revenue: 3–5%",
     GOLD_ACCENT),
]

for i, (title, qty, desc, col) in enumerate(byproducts):
    x = Cm(2) + Cm(i * 7.8)
    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, Cm(5), Cm(7.2), Cm(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Quantity
    add_textbox(slide, x, Cm(6.7), Cm(7.2), Cm(0.8),
                qty, font_size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Cm(0.3), Cm(7.5), Cm(7), Cm(5),
                desc, font_size=9, color=INDUSTRIAL_GRAY)

# Total revenue note
add_textbox(slide, Cm(2), Cm(13.5), Cm(29), Cm(1.5),
            "Total kontribusi by-product: 15–25% dari revenue pabrik gula. "
            "Ini menjadikan pabrik gula sebagai industri circular economy — hampir tidak ada waste.",
            font_size=10, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)

add_key_takeaway(slide, "By-product bukan limbah, tapi sumber pendapatan tambahan yang signifikan (15–25% revenue).")



# ============================================================
# SLIDE 14: GMP (GOOD MANUFACTURING PRACTICE)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Good Manufacturing Practice (GMP)",
                   "Standar Produksi — Menjamin mutu dan keamanan pangan")
add_footer(slide, 14)

gmp_lines = [
    "Apa itu GMP?",
    "Pedoman cara produksi yang baik untuk memastikan produk aman,",
    "bermutu, dan konsisten — wajib untuk industri pangan (termasuk gula).",
    "",
    "5 Pilar GMP di Pabrik Gula:",
]
add_multi_text(slide, Cm(2), Cm(5), Cm(29), Cm(4),
               gmp_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

pillars = [
    ("1. Bangunan & Fasilitas",
     "• Desain higienis, mudah dibersihkan\n• Alur proses searah (no cross-contamination)\n• Ventilasi & penerangan memadai"),
    ("2. Peralatan & Mesin",
     "• Food-grade material (stainless steel)\n• Mudah dibersihkan & dirawat\n• Kalibrasi rutin alat ukur"),
    ("3. Personel / SDM",
     "• Pelatihan hygiene & GMP rutin\n• APD lengkap di area produksi\n• Pemeriksaan kesehatan berkala"),
    ("4. Proses Produksi",
     "• SOP tertulis & terdokumentasi\n• Kontrol titik kritis (CCP)\n• Traceability batch produksi"),
    ("5. Sanitasi & Pengendalian Hama",
     "• Jadwal pembersihan teratur\n• Pest control terintegrasi\n• Pengelolaan limbah cair & padat"),
]

for i, (title, desc) in enumerate(pillars):
    x = Cm(2) + Cm(i * 6.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Cm(9), Cm(5.8), Cm(6.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = SUGARCANE_GREEN
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    tf.margin_left = Pt(6)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(8)
    p2.font.color.rgb = INDUSTRIAL_GRAY
    p2.space_before = Pt(4)

add_key_takeaway(slide, "GMP memastikan setiap tahap produksi menghasilkan gula yang aman & konsisten — dari bahan baku hingga produk akhir.", Cm(16))



# ============================================================
# SLIDE 15: QUALITY CONTROL
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Quality Control (Pengendalian Mutu)",
                   "Laboratorium & Monitoring — Menjaga standar di setiap stasiun")
add_footer(slide, 15)

qc_lines = [
    "Apa yang terjadi di proses ini?",
    "",
    "Tim QC / laboratorium pabrik memantau kualitas di SETIAP stasiun proses,",
    "dari tebu masuk hingga gula siap kirim.",
    "",
    "Parameter yang Diuji di Setiap Titik:",
]
add_multi_text(slide, Cm(2), Cm(5), Cm(29), Cm(4),
               qc_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

# QC table-like layout
qc_points = [
    ("Bahan Baku (Tebu)", "Pol, Brix, Trash %, Fiber %"),
    ("Nira Mentah", "Brix, Pol, pH, Purity"),
    ("Nira Jernih", "Turbidity, Color, pH"),
    ("Nira Kental", "Brix (60-65%), Purity"),
    ("Massecuite", "Brix, Crystal size, Purity"),
    ("Gula Produk", "Pol ≥99,3%, ICUMSA, Moisture <0,05%"),
    ("Molasses", "Brix, Pol, Target Purity <35%"),
    ("Air Boiler", "Hardness, pH, TDS, Silica"),
]

for i, (point, params) in enumerate(qc_points):
    row = i // 4
    col = i % 4
    x = Cm(2) + Cm(col * 7.8)
    y = Cm(9.5) + Cm(row * 3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Cm(7.2), Cm(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = SUGARCANE_GREEN
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(6)
    p = tf.paragraphs[0]
    p.text = point
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p2 = tf.add_paragraph()
    p2.text = params
    p2.font.size = Pt(8)
    p2.font.color.rgb = INDUSTRIAL_GRAY
    p2.space_before = Pt(3)

add_key_takeaway(slide, "QC berjalan 24 jam — setiap 2 jam sampling dilakukan. Data lab = dasar keputusan operasional.", Cm(16))



# ============================================================
# SLIDE 16: PENGELOLAAN LIMBAH
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Pengelolaan Limbah",
                   "Waste Management — Mengelola dampak lingkungan")
add_footer(slide, 16)

waste_lines = [
    "Jenis Limbah Pabrik Gula & Pengelolaannya:",
]
add_multi_text(slide, Cm(2), Cm(5), Cm(29), Cm(1.5),
               waste_lines, font_size=11, bold_first=True, color=DARK_GREEN)

waste_types = [
    ("Limbah Cair", SUGARCANE_GREEN, [
        "Sumber: air pencucian, kondensor, boiler blowdown",
        "Penanganan:",
        "• IPAL (Instalasi Pengolahan Air Limbah)",
        "• Sistem: screening → equalisasi → aerasi → clarifier",
        "• BOD/COD diturunkan sebelum dibuang",
        "• Target: sesuai baku mutu Permen LH",
        "• Bisa digunakan ulang untuk irigasi kebun"
    ]),
    ("Limbah Padat", WARM_BROWN, [
        "Sumber: blotong, abu boiler, trash tebu",
        "Penanganan:",
        "• Blotong → pupuk organik / kompos",
        "• Abu boiler (fly ash) → campuran pupuk / landfill",
        "• Trash → dikembalikan ke lahan sebagai mulsa",
        "• Zero waste to landfill (target modern)",
        "• Bagasse habis dipakai sebagai bahan bakar"
    ]),
    ("Emisi Gas (Udara)", INDUSTRIAL_GRAY, [
        "Sumber: cerobong boiler, debu pengering",
        "Penanganan:",
        "• Wet scrubber / cyclone pada cerobong",
        "• Electrostatic precipitator (ESP)",
        "• Monitoring emisi berkala (PM, SO₂, NOx)",
        "• Target: di bawah baku mutu emisi",
        "• Bag filter pada area drying & packaging"
    ]),
]

for i, (title, col, items) in enumerate(waste_types):
    x = Cm(2) + Cm(i * 10.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Cm(6.5), Cm(10), Cm(9.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = col
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(10)
    tf.margin_left = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = col
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = item
        p2.font.size = Pt(9)
        p2.font.color.rgb = INDUSTRIAL_GRAY
        p2.space_before = Pt(2)

add_key_takeaway(slide, "Pabrik gula modern menerapkan prinsip zero waste — semua limbah diolah menjadi produk bernilai atau dibuang sesuai standar.", Cm(16.5))



# ============================================================
# SLIDE 17: SUSTAINABILITY & EFISIENSI
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Sustainability & Efisiensi Energi",
                   "Keberlanjutan — Menyeimbangkan produksi dan lingkungan")
add_footer(slide, 17)

sustain_lines = [
    "Mengapa Sustainability Penting di Pabrik Gula?",
    "",
    "Pabrik gula memiliki potensi besar untuk menjadi industri berkelanjutan",
    "karena sifat alamiahnya yang circular (produk samping = bahan baku kembali).",
    "",
    "Aspek Sustainability di Pabrik Gula:",
]
add_multi_text(slide, Cm(2), Cm(5), Cm(29), Cm(4),
               sustain_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

aspects = [
    ("Energi Terbarukan", [
        "• Bagasse sebagai biomass fuel (carbon neutral)",
        "• Co-generation: listrik dari uap boiler",
        "• Potensi solar panel di lahan idle",
        "• Target: energy self-sufficient"
    ]),
    ("Efisiensi Air", [
        "• Recycle condensate untuk boiler feed water",
        "• Reuse air proses untuk irigasi",
        "• Minimasi konsumsi air bersih",
        "• Water footprint tracking"
    ]),
    ("Circular Economy", [
        "• Blotong → pupuk → kembali ke kebun",
        "• Molasses → industri turunan",
        "• Abu → soil amendment",
        "• Hampir zero waste"
    ]),
    ("Carbon Footprint", [
        "• Bagasse = biomass (CO₂ netral)",
        "• Mengurangi ketergantungan fossil fuel",
        "• PROPER KLHK (Blue/Green rating)",
        "• Potensi carbon credit trading"
    ]),
]

for i, (title, items) in enumerate(aspects):
    x = Cm(2) + Cm(i * 7.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Cm(9.5), Cm(7.2), Cm(5.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = SUGARCANE_GREEN
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    tf.margin_left = Pt(6)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = item
        p2.font.size = Pt(9)
        p2.font.color.rgb = INDUSTRIAL_GRAY
        p2.space_before = Pt(2)

add_key_takeaway(slide, "Pabrik gula punya DNA sustainability — bagasse sebagai energi + by-product sebagai revenue = model circular economy.", Cm(15.8))



# ============================================================
# SLIDE 18: KESELAMATAN KERJA (K3)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Keselamatan & Kesehatan Kerja (K3)",
                   "Safety First — Operasi aman, produktivitas terjaga")
add_footer(slide, 18)

k3_lines = [
    "Potensi Bahaya di Pabrik Gula:",
    "",
    "• Mesin berputar (gilingan, centrifuge, turbin) → risiko terjerat",
    "• Suhu tinggi (boiler, evaporator, pan masakan) → risiko luka bakar",
    "• Ketinggian (tangki, cerobong, conveyor) → risiko jatuh",
    "• Kebisingan (area gilingan >85 dB) → gangguan pendengaran",
    "• Bahan kimia (kapur, belerang, asam) → iritasi kulit/mata",
    "• Debu gula (area dryer, packaging) → risiko respirasi",
]
add_multi_text(slide, Cm(2), Cm(5), Cm(15), Cm(7),
               k3_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

# Right side: Safety measures
safety_lines = [
    "Penerapan K3 di Pabrik Gula:",
    "",
    "1. APD Wajib:",
    "   Helm, safety shoes, earplug, goggle, sarung tangan",
    "",
    "2. Sistem Izin Kerja:",
    "   Hot work permit, confined space permit, LOTO",
    "",
    "3. Training & Drill:",
    "   Induction safety, fire drill, emergency response",
    "",
    "4. Inspeksi Rutin:",
    "   Safety patrol harian, audit K3 bulanan",
    "",
    "5. Pelaporan & Investigasi:",
    "   Near-miss reporting, incident investigation",
    "",
    "6. Safety Committee:",
    "   P2K3 (Panitia Pembina K3) aktif & meeting rutin",
]
add_multi_text(slide, Cm(17), Cm(5), Cm(15), Cm(11),
               safety_lines, font_size=10, color=INDUSTRIAL_GRAY, bold_first=True)

add_key_takeaway(slide, "Zero accident adalah target — K3 bukan sekadar aturan, tapi budaya kerja yang melindungi semua karyawan.", Cm(16.5))



# ============================================================
# SLIDE 19: TANTANGAN OPERASIONAL
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Tantangan Operasional Pabrik Gula",
                   "Hambatan & Bottleneck — Yang harus dihadapi industri")
add_footer(slide, 19)

challenges = [
    ("Produktivitas Rendah",
     "Rendemen RI 6–8% vs Thailand 10–12%. Penyebab: varietas tebu, umur ratoon, masa tebang.",
     SUGARCANE_GREEN),
    ("Teknologi Tua",
     "Banyak mesin berusia >30 tahun. Breakdown tinggi, efisiensi rendah, spare parts langka.",
     WARM_BROWN),
    ("Kualitas Bahan Baku",
     "Tebu kotor (trash >5%), pol rendah, terlambat giling. Petani belum optimal menerapkan standar.",
     DARK_GREEN),
    ("Margin Tipis",
     "HPP Rp 9–11rb/kg vs harga jual terkontrol. Sensitif terhadap rendemen & harga tebu.",
     GOLD_ACCENT),
    ("SDM & Regenerasi",
     "Operator senior pensiun, generasi muda kurang tertarik. Perlu training intensif & sertifikasi.",
     INDUSTRIAL_GRAY),
    ("Risiko Iklim",
     "El Niño/La Niña ganggu jadwal giling. Hujan saat musim giling = tebu basah, rendemen turun.",
     SUGARCANE_GREEN),
]

for i, (title, desc, col) in enumerate(challenges):
    row = i // 3
    c = i % 3
    x = Cm(2) + Cm(c * 10.5)
    y = Cm(5) + Cm(row * 5.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Cm(9.8), Cm(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = col
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(10)
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = col
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = INDUSTRIAL_GRAY
    p2.space_before = Pt(6)

add_key_takeaway(slide, "Tantangan terbesar: rendemen rendah & teknologi tua. Solusi: modernisasi + digitalisasi + kualitas tebu.", Cm(16.5))



# ============================================================
# SLIDE 20: KESIMPULAN
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_section_header(slide, "Kesimpulan & Ringkasan",
                   "5 Key Takeaways dari presentasi ini")
add_footer(slide, 20)

conclusions = [
    ("1", "Proses Terintegrasi",
     "Pabrik gula menjalankan 9 tahapan proses kontinu selama 24/7 di musim giling. Dari tebu masuk hingga gula terkemas = 24–48 jam."),
    ("2", "Multi Revenue Stream",
     "Gula (70–80% revenue) + by-product 15–25% (molasses, bagasse, blotong, listrik). Model circular economy."),
    ("3", "KPI Raja: Rendemen",
     "Rendemen 8–10% = ukuran keberhasilan utama. Naik 1% rendemen bisa naikkan profit secara signifikan."),
    ("4", "GMP & QC Ketat",
     "Kualitas dijaga dari bahan baku hingga produk akhir. Lab beroperasi 24 jam, sampling tiap 2 jam."),
    ("5", "Peluang Besar",
     "Digitalisasi, modernisasi, co-generation, dan diversifikasi = potensi peningkatan EBITDA 20–30%."),
]

for i, (num, title, desc) in enumerate(conclusions):
    y = Cm(5) + Cm(i * 2.4)
    # Number
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2.5), y, Cm(1.5), Cm(1.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = SUGARCANE_GREEN
    circ.line.fill.background()
    tf = circ.text_frame
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_textbox(slide, Cm(4.5), y, Cm(7), Cm(1.5),
                title, font_size=12, bold=True, color=DARK_GREEN)
    # Description
    add_textbox(slide, Cm(11), y, Cm(20), Cm(2),
                desc, font_size=10, color=INDUSTRIAL_GRAY)



# ============================================================
# SLIDE 21: CLOSING SLIDE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_GREEN)

add_textbox(slide, Cm(3), Cm(4), Cm(28), Cm(2),
            "Terima Kasih", font_size=36, bold=True,
            color=WHITE, align=PP_ALIGN.LEFT)

add_textbox(slide, Cm(3), Cm(7), Cm(28), Cm(2),
            "Siap menerima pertanyaan, diskusi, & masukan.",
            font_size=16, color=RGBColor(0xA5, 0xD6, 0xA7), align=PP_ALIGN.LEFT)

# Decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(10), Cm(8), Cm(0.15))
line.fill.solid()
line.fill.fore_color.rgb = GOLD_ACCENT
line.line.fill.background()

closing_lines = [
    "Referensi Utama:",
    "• Hugot, E. — Handbook of Cane Sugar Engineering (1960)",
    "• Good Management Practices Manual for the Cane Sugar Industry",
    "• Laporan Kerja Praktik: Rika Aisyah & Vera Arum (2031910045/052)",
    "• Denik Sentoso — Naskah Publikasi Jurnal",
    "• Jurnal Teknologi Gula — Article Text 6234",
    "",
    "Presentasi Internal  •  2026  •  Overview Operasional & Proses Produksi Pabrik Gula"
]
add_multi_text(slide, Cm(3), Cm(11), Cm(28), Cm(6),
               closing_lines, font_size=10, color=RGBColor(0xC8, 0xE6, 0xC9),
               bold_first=True)


# ============================================================
# SAVE PRESENTATION
# ============================================================
output_path = "buat-tugas/Overview_Operasional_dan_Proses_Produksi_Pabrik_Gula.pptx"
prs.save(output_path)
print(f"Presentation saved successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")
