"""
Build: Overview_Proses_Alur_Bisnis_Pabrik_Gula.pptx
Professional 16-slide deck - Bahasa Indonesia
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# === COLOR PALETTE ===
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
MID_GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
SUGARCANE = RGBColor(0x7C, 0xB3, 0x42)
GOLD = RGBColor(0xC6, 0x8A, 0x02)
BROWN = RGBColor(0x6D, 0x4C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF1, 0xF8, 0xE9)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)

FONT_TITLE = 'Calibri'
FONT_BODY = 'Calibri'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# === HELPER FUNCTIONS ===
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_box(slide, left, top, width, height, fill_color, text='',
                  font_size=Pt(11), font_color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT_BODY
        p.alignment = align
    return shape

def add_title_bar(slide, title_text, subtitle_text=''):
    """Add a green title bar at top of slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_GREEN
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_top = Pt(16)
    tf.margin_left = Pt(40)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_TITLE
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xC8, 0xE6, 0xC9)
        p2.font.name = FONT_BODY

def add_body_text(slide, left, top, width, height, bullets, font_size=Pt(14)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = font_size
        p.font.color.rgb = DARK_TEXT
        p.font.name = FONT_BODY
        p.space_after = Pt(8)
    return txBox

def add_key_takeaway(slide, left, top, width, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = "💡 Key Takeaway: " + text
    p.font.size = Pt(12)
    p.font.color.rgb = BROWN
    p.font.bold = True
    p.font.name = FONT_BODY
    return shape


def add_flow_arrow(slide, left, top):
    """Small right arrow between flow boxes."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.4), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()
    arrow.shadow.inherit = False

def add_down_arrow(slide, left, top):
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.3), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()
    arrow.shadow.inherit = False

# ============================================================
# SLIDE 1: COVER
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, DARK_GREEN)

# Accent bar
accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(2.6))
accent.fill.solid()
accent.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0x0A)
accent.line.fill.background()
accent.shadow.inherit = False

# Title
txBox = slide1.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "OVERVIEW PROSES ALUR BISNIS"
p.font.size = Pt(38)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = FONT_TITLE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "PABRIK GULA"
p2.font.size = Pt(44)
p2.font.color.rgb = GOLD
p2.font.bold = True
p2.font.name = FONT_TITLE
p2.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = slide1.shapes.add_textbox(Inches(2), Inches(5.2), Inches(9), Inches(0.8))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "Dari Kebun Tebu hingga Distribusi Gula — Panduan Operasional untuk Pemula"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0xA5, 0xD6, 0xA7)
p3.font.name = FONT_BODY
p3.alignment = PP_ALIGN.CENTER

# Footer
txBox3 = slide1.shapes.add_textbox(Inches(3), Inches(6.5), Inches(7), Inches(0.5))
tf3 = txBox3.text_frame
p4 = tf3.paragraphs[0]
p4.text = "Presentasi Kerja Praktik | 2024"
p4.font.size = Pt(12)
p4.font.color.rgb = RGBColor(0x81, 0xC7, 0x84)
p4.font.name = FONT_BODY
p4.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, WHITE)
add_title_bar(slide2, "DAFTAR ISI", "Agenda Presentasi")

toc_items = [
    "1. Pengantar Industri Gula",
    "2. Alur Bisnis End-to-End",
    "3. Penerimaan Bahan Baku & Tebu",
    "4. Persiapan Tebu & Penggilingan",
    "5. Pemurnian Nira (Clarification)",
    "6. Penguapan (Evaporation)",
    "7. Kristalisasi & Vacuum Pan",
    "8. Sentrifugasi & Pemisahan Gula",
    "9. Pengeringan, Pengemasan & Penyimpanan",
    "10. Utilitas & Sistem Pendukung",
    "11. Alur Bisnis & Operasional",
    "12. Tantangan Operasional Pabrik Gula",
    "13. Rangkuman & Key Takeaways",
]

# Two columns
col1 = toc_items[:7]
col2 = toc_items[7:]
add_body_text(slide2, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), col1, Pt(15))
add_body_text(slide2, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5), col2, Pt(15))

# Decorative side bar
side = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.8), Inches(1.3), Inches(0.15), Inches(6.2))
side.fill.solid()
side.fill.fore_color.rgb = SUGARCANE
side.line.fill.background()
side.shadow.inherit = False

# ============================================================
# SLIDE 3: INTRODUCTION TO SUGAR INDUSTRY
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, WHITE)
add_title_bar(slide3, "PENGANTAR INDUSTRI GULA", "Apa itu Pabrik Gula & Perannya dalam Supply Chain")

bullets3 = [
    "▪ Pabrik gula = industri terintegrasi: dari kebun tebu → proses → distribusi",
    "▪ Output utama: Gula Kristal Putih (SHS) dengan pol ≥99,3%",
    "▪ 4 produk samping: molasses, bagasse, blotong, listrik (15–25% revenue)",
    "",
    "▪ Produksi nasional: ~2,2 juta ton/tahun",
    "▪ Konsumsi nasional: ~5–6 juta ton/tahun",
    "▪ Impor: ~50% dari kebutuhan",
    "▪ Pabrik aktif di Indonesia: ~60 unit",
    "",
    "▪ Musim giling: April – November (6–7 bulan)",
    "▪ Off-season: maintenance & overhaul",
]
add_body_text(slide3, Inches(0.8), Inches(1.6), Inches(7), Inches(5.2), bullets3, Pt(14))

# Stats boxes on right
add_shape_box(slide3, Inches(8.5), Inches(1.8), Inches(3.8), Inches(1),
              MID_GREEN, "Rendemen Target\n8–10%", Pt(16), WHITE, True)
add_shape_box(slide3, Inches(8.5), Inches(3.0), Inches(3.8), Inches(1),
              GOLD, "Kapasitas Giling\n3.000–6.000 TCD", Pt(16), WHITE, True)
add_shape_box(slide3, Inches(8.5), Inches(4.2), Inches(3.8), Inches(1),
              BROWN, "Lead Time\n24–48 jam", Pt(16), WHITE, True)
add_shape_box(slide3, Inches(8.5), Inches(5.4), Inches(3.8), Inches(1),
              DARK_GREEN, "Margin EBITDA\n15–25%", Pt(16), WHITE, True)


# ============================================================
# SLIDE 4: END-TO-END BUSINESS PROCESS OVERVIEW
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, WHITE)
add_title_bar(slide4, "ALUR BISNIS END-TO-END", "High-Level Process Flow — Dari Kebun hingga Konsumen")

# Flow diagram - 2 rows of boxes with arrows
flow_top = [
    ("Penanaman\nTebu", "10-14 bln"),
    ("Tebang &\nAngkut", "<24 jam"),
    ("Penerimaan\n& Timbang", "Cek Pol/Brix"),
    ("Penggilingan\n(Milling)", "Ekstraksi Nira"),
    ("Pemurnian\n(Clarification)", "Hilangkan Impuritas"),
]
flow_bot = [
    ("Penguapan\n(Evaporation)", "Brix 12%→65%"),
    ("Kristalisasi\n(Vacuum Pan)", "Bentuk Kristal"),
    ("Sentrifugasi", "Pisah Kristal\n& Molasses"),
    ("Pengeringan\n& Packing", "Kadar Air <0,5%"),
    ("Distribusi\n& Penjualan", "Gudang → Pasar"),
]

box_w = Inches(2.1)
box_h = Inches(1.1)
start_x = Inches(0.5)
y1 = Inches(1.8)
y2 = Inches(4.5)
gap = Inches(0.5)

for i, (label, sub) in enumerate(flow_top):
    x = start_x + i * (box_w + gap)
    color = MID_GREEN if i % 2 == 0 else DARK_GREEN
    shape = add_shape_box(slide4, x, y1, box_w, box_h, color, label, Pt(12), WHITE, True)
    # sub-label
    sub_box = slide4.shapes.add_textbox(x, y1 + box_h + Pt(4), box_w, Inches(0.4))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = sub
    sp.font.size = Pt(10)
    sp.font.color.rgb = GRAY
    sp.font.name = FONT_BODY
    sp.alignment = PP_ALIGN.CENTER
    if i < len(flow_top) - 1:
        add_flow_arrow(slide4, x + box_w + Inches(0.05), y1 + Inches(0.4))

for i, (label, sub) in enumerate(flow_bot):
    x = start_x + i * (box_w + gap)
    color = DARK_GREEN if i % 2 == 0 else MID_GREEN
    shape = add_shape_box(slide4, x, y2, box_w, box_h, color, label, Pt(12), WHITE, True)
    sub_box = slide4.shapes.add_textbox(x, y2 + box_h + Pt(4), box_w, Inches(0.4))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = sub
    sp.font.size = Pt(10)
    sp.font.color.rgb = GRAY
    sp.font.name = FONT_BODY
    sp.alignment = PP_ALIGN.CENTER
    if i < len(flow_bot) - 1:
        add_flow_arrow(slide4, x + box_w + Inches(0.05), y2 + Inches(0.4))

# Connecting arrow between rows
add_down_arrow(slide4, Inches(6.3), Inches(3.3))

add_key_takeaway(slide4, Inches(0.5), Inches(6.4), Inches(12),
                 "Total lead time tebu tebang → gula jadi = 24–48 jam. Setiap jam delay = pol turun 0,1%")


# ============================================================
# SLIDE 5: RAW MATERIAL SUPPLY & CANE RECEIVING
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, WHITE)
add_title_bar(slide5, "PENERIMAAN BAHAN BAKU & TEBU",
              "Farmer Supply → Transportasi → Unloading → Cane Yard")

bullets5 = [
    "▪ Tebu ditanam 10–14 bulan sebelum dipanen (ratoon/plant cane)",
    "▪ Petani/supplier mengirim tebu via truk ke pabrik",
    "▪ Tebu harus digiling <24 jam setelah tebang (cegah pol turun)",
    "",
    "▪ Proses Penerimaan:",
    "   1. Penimbangan truk (brutto & tarra)",
    "   2. Sampling & analisis (pol, brix, trash %)",
    "   3. Unloading dengan cane tipper / crane",
    "   4. Penyimpanan sementara di cane yard",
    "",
    "▪ Cane yard mengatur antrian FIFO (first-in first-out)",
    "▪ Target: antrian maksimal 12–24 jam",
]
add_body_text(slide5, Inches(0.8), Inches(1.6), Inches(7), Inches(5), bullets5, Pt(13))

# Mini flow on right
flow_labels = ["Petani / Kebun", "Transportasi Truk", "Timbangan & QC",
               "Unloading (Tipper)", "Cane Yard / Lori"]
fy = Inches(1.8)
for i, lbl in enumerate(flow_labels):
    add_shape_box(slide5, Inches(9), fy + i * Inches(1.05), Inches(3.5), Inches(0.7),
                  MID_GREEN if i % 2 == 0 else SUGARCANE, lbl, Pt(12), WHITE, True)
    if i < len(flow_labels) - 1:
        add_down_arrow(slide5, Inches(10.6), fy + i * Inches(1.05) + Inches(0.72))

add_key_takeaway(slide5, Inches(0.8), Inches(6.4), Inches(8),
                 "Kualitas tebu saat diterima menentukan rendemen akhir pabrik.")

# ============================================================
# SLIDE 6: CANE PREPARATION & MILLING
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, WHITE)
add_title_bar(slide6, "PERSIAPAN TEBU & PENGGILINGAN (MILLING)",
              "Cane Carrier → Knives/Shredder → Milling Tandem → Nira Mentah")

bullets6_left = [
    "▪ Cane Carrier: konveyor membawa tebu ke stasiun gilingan",
    "▪ Cane Knife (pisau tebu): memotong tebu jadi potongan kecil",
    "▪ Shredder/Unigrator: mencacah serat tebu agar sel terbuka",
    "",
    "▪ Milling Tandem (4–6 unit gilingan):",
    "   – Setiap unit = 3 roll (top, feed, discharge)",
    "   – Imbibisi air panas untuk maksimalkan ekstraksi",
    "   – Hasil: Nira Mentah (mixed juice) + Ampas (bagasse)",
    "",
    "▪ Target Recovery Pol: 92–96%",
    "▪ Kapasitas: 3.000–6.000 TCD (Ton Cane per Day)",
]
add_body_text(slide6, Inches(0.8), Inches(1.6), Inches(7), Inches(5), bullets6_left, Pt(13))

# Milling process mini-flow
mill_steps = ["Cane Carrier", "Cane Knife", "Shredder",
              "Gilingan 1-2-3-4", "Nira Mentah\n+ Bagasse"]
my = Inches(1.8)
for i, lbl in enumerate(mill_steps):
    add_shape_box(slide6, Inches(9), my + i * Inches(1.05), Inches(3.5), Inches(0.7),
                  DARK_GREEN if i % 2 == 0 else MID_GREEN, lbl, Pt(12), WHITE, True)
    if i < len(mill_steps) - 1:
        add_down_arrow(slide6, Inches(10.6), my + i * Inches(1.05) + Inches(0.72))

add_key_takeaway(slide6, Inches(0.8), Inches(6.4), Inches(8),
                 "100 ton tebu → ~30 ton bagasse (bahan bakar) + nira untuk diproses lebih lanjut")


# ============================================================
# SLIDE 7: JUICE TREATMENT & CLARIFICATION
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, WHITE)
add_title_bar(slide7, "PEMURNIAN NIRA (CLARIFICATION)",
              "Penghilangan Impuritas → Pemanasan → Nira Jernih")

bullets7 = [
    "▪ Nira mentah mengandung gula + kotoran (tanah, lilin, protein)",
    "▪ Tujuan: menghasilkan nira jernih (clear juice) untuk diuapkan",
    "",
    "▪ Proses Utama:",
    "   1. Pre-liming: penambahan susu kapur Ca(OH)₂ → pH ~7",
    "   2. Sulfitasi (opsional): SO₂ untuk pemucatan warna",
    "   3. Pemanasan (juice heater): 100–105°C",
    "   4. Flash tank: lepas udara terlarut",
    "   5. Clarifier (pengendapan): kotoran mengendap",
    "   6. Rotary Vacuum Filter: saring endapan → blotong",
    "",
    "▪ Output: Nira Jernih (clear juice) → ke evaporator",
    "▪ By-product: Blotong (filter cake) → pupuk organik",
]
add_body_text(slide7, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.2), bullets7, Pt(13))

# Diagram
clar_steps = ["Nira Mentah", "Liming + Sulfitasi", "Juice Heater\n(100-105°C)",
              "Clarifier", "Nira Jernih →\nEvaporator"]
cy = Inches(1.8)
for i, lbl in enumerate(clar_steps):
    add_shape_box(slide7, Inches(9.2), cy + i * Inches(1.05), Inches(3.3), Inches(0.7),
                  MID_GREEN if i % 2 == 0 else SUGARCANE, lbl, Pt(11), WHITE, True)
    if i < len(clar_steps) - 1:
        add_down_arrow(slide7, Inches(10.7), cy + i * Inches(1.05) + Inches(0.72))

add_key_takeaway(slide7, Inches(0.8), Inches(6.5), Inches(8),
                 "Blotong bukan limbah — dijual sebagai pupuk organik (2–3% revenue tambahan)")

# ============================================================
# SLIDE 8: EVAPORATION PROCESS
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, WHITE)
add_title_bar(slide8, "PENGUAPAN (EVAPORATION)",
              "Konsentrasi Nira — Multiple Effect Evaporator")

bullets8 = [
    "▪ Tujuan: menguapkan air dari nira jernih → nira kental (thick juice/syrup)",
    "▪ Brix naik dari ~12% → ~65% (konsentrasi gula meningkat)",
    "",
    "▪ Sistem: Multiple Effect Evaporator (MEE)",
    "   – Biasanya 4–5 badan (effect/vessel)",
    "   – Uap dari badan 1 digunakan memanaskan badan 2, dst.",
    "   – Efisiensi energi: 1 kg uap bisa menguapkan 4–5 kg air",
    "",
    "▪ Prinsip kerja:",
    "   – Tekanan menurun tiap badan (vakum progresif)",
    "   – Titik didih turun → penguapan tetap terjadi",
    "   – Badan terakhir beroperasi di bawah vakum",
    "",
    "▪ Output: Nira Kental (syrup) siap masuk vacuum pan",
]
add_body_text(slide8, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.2), bullets8, Pt(13))

# MEE diagram (horizontal flow)
evap_labels = ["Effect 1\n(Tekanan\nTinggi)", "Effect 2", "Effect 3",
               "Effect 4\n(Vakum)", "Nira\nKental"]
ex = Inches(0.6)
ey = Inches(6.0)
for i, lbl in enumerate(evap_labels):
    color = DARK_GREEN if i < 4 else GOLD
    add_shape_box(slide8, ex + i * Inches(2.6), ey, Inches(2.2), Inches(0.9),
                  color, lbl, Pt(10), WHITE, True)
    if i < len(evap_labels) - 1:
        add_flow_arrow(slide8, ex + i * Inches(2.6) + Inches(2.25), ey + Inches(0.3))


# ============================================================
# SLIDE 9: CRYSTALLIZATION & VACUUM PAN
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, WHITE)
add_title_bar(slide9, "KRISTALISASI & VACUUM PAN",
              "Pembentukan Kristal Gula — Massecuite Process")

bullets9 = [
    "▪ Tujuan: membentuk kristal gula dari nira kental",
    "▪ Dilakukan dalam Vacuum Pan (bejana vakum)",
    "",
    "▪ Proses:",
    "   1. Nira kental dimasukkan ke vacuum pan",
    "   2. Pemanasan vakum → penguapan air lanjut",
    "   3. Penambahan bibit kristal (seed/fondan)",
    "   4. Kristal tumbuh perlahan dalam larutan induk",
    "   5. Hasil: Massecuite (campuran kristal + mother liquor)",
    "",
    "▪ Sistem Masakan (bertingkat):",
    "   – Masakan A: kristal terbesar, kualitas terbaik",
    "   – Masakan B: kristal sedang",
    "   – Masakan C/D: kristal kecil (diputar ulang)",
    "",
    "▪ Massecuite A → sentrifugasi → Gula SHS",
]
add_body_text(slide9, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.2), bullets9, Pt(13))

# Crystal hierarchy
cryst_labels = ["Nira Kental\n(Syrup)", "Vacuum Pan\n(Masakan A/B/C)", "Massecuite",
                "→ Sentrifugasi"]
cry = Inches(2.0)
for i, lbl in enumerate(cryst_labels):
    add_shape_box(slide9, Inches(9.2), cry + i * Inches(1.2), Inches(3.3), Inches(0.8),
                  DARK_GREEN if i % 2 == 0 else MID_GREEN, lbl, Pt(12), WHITE, True)
    if i < len(cryst_labels) - 1:
        add_down_arrow(slide9, Inches(10.7), cry + i * Inches(1.2) + Inches(0.82))

add_key_takeaway(slide9, Inches(0.8), Inches(6.5), Inches(8),
                 "Kristalisasi bertingkat memaksimalkan recovery — gula yang tidak terkristal diproses ulang")

# ============================================================
# SLIDE 10: CENTRIFUGATION & SUGAR SEPARATION
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, WHITE)
add_title_bar(slide10, "SENTRIFUGASI & PEMISAHAN GULA",
              "Memisahkan Kristal Gula dari Molasses")

bullets10 = [
    "▪ Tujuan: memisahkan kristal gula murni dari larutan induk (molasses)",
    "▪ Alat: Centrifugal / Putaran (basket type, batch/continuous)",
    "",
    "▪ Proses:",
    "   1. Massecuite dimasukkan ke basket centrifugal",
    "   2. Putaran tinggi (1.000–1.500 rpm)",
    "   3. Gaya sentrifugal mendorong molasses keluar melalui saringan",
    "   4. Kristal gula tertahan di dalam basket",
    "   5. Pencucian (steam/water wash) untuk poles kristal",
    "",
    "▪ Hasil:",
    "   – Kristal gula basah → ke pengering",
    "   – Molasses A → masakan B (recycle)",
    "   – Final Molasses (tetes akhir) → dijual (pakan, etanol)",
    "",
    "▪ 100 ton tebu → ~4–5 ton molasses",
]
add_body_text(slide10, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.2), bullets10, Pt(13))

# Simple diagram
sep_items = [("Massecuite", MID_GREEN), ("Centrifugal\n(1000-1500 rpm)", DARK_GREEN),
             ("Kristal Gula\n(basah)", SUGARCANE), ("Molasses\n(tetes)", GOLD)]
sy = Inches(2.0)
for i, (lbl, col) in enumerate(sep_items):
    add_shape_box(slide10, Inches(9.2), sy + i * Inches(1.2), Inches(3.3), Inches(0.8),
                  col, lbl, Pt(12), WHITE, True)
    if i < len(sep_items) - 1:
        add_down_arrow(slide10, Inches(10.7), sy + i * Inches(1.2) + Inches(0.82))


# ============================================================
# SLIDE 11: DRYING, PACKAGING, AND STORAGE
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11, WHITE)
add_title_bar(slide11, "PENGERINGAN, PENGEMASAN & PENYIMPANAN",
              "Finishing Process — Sugar Drying, Packing, Warehouse")

bullets11 = [
    "▪ Pengeringan (Sugar Dryer):",
    "   – Kristal gula basah dilewatkan rotary dryer",
    "   – Udara panas mengurangi kadar air hingga <0,5%",
    "   – Dilanjutkan cooler → suhu turun ke ambient",
    "",
    "▪ Grading & Screening:",
    "   – Ayakan memisahkan ukuran kristal (halus/kasar)",
    "   – Gula off-spec dilebur ulang (re-melt)",
    "",
    "▪ Pengemasan:",
    "   – Standar: karung 50 kg",
    "   – Automatic weighing & sewing machine",
    "   – Labeling: batch, tanggal produksi, pol, ICUMSA",
    "",
    "▪ Gudang Penyimpanan:",
    "   – Suhu & kelembaban terkontrol",
    "   – FIFO system, pallet stacking",
    "   – Siap distribusi ke pasar/industri",
]
add_body_text(slide11, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.5), bullets11, Pt(13))

# Flow on right
dry_steps = ["Kristal Basah", "Rotary Dryer\n& Cooler", "Grading\n(Screening)",
             "Pengemasan\n(50 kg)", "Gudang\n(Distribusi)"]
dy = Inches(1.8)
for i, lbl in enumerate(dry_steps):
    add_shape_box(slide11, Inches(9.2), dy + i * Inches(1.05), Inches(3.3), Inches(0.7),
                  DARK_GREEN if i % 2 == 0 else SUGARCANE, lbl, Pt(11), WHITE, True)
    if i < len(dry_steps) - 1:
        add_down_arrow(slide11, Inches(10.7), dy + i * Inches(1.05) + Inches(0.72))

# ============================================================
# SLIDE 12: UTILITY & SUPPORTING SYSTEMS
# ============================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide12, WHITE)
add_title_bar(slide12, "UTILITAS & SISTEM PENDUKUNG",
              "Steam, Boiler, Turbine, Bagasse & Energy Cycle")

bullets12 = [
    "▪ Pabrik gula hampir self-sufficient untuk energi!",
    "",
    "▪ Siklus Energi (Cogeneration):",
    "   1. Bagasse (ampas tebu) → dibakar di Boiler",
    "   2. Boiler → menghasilkan uap bertekanan tinggi",
    "   3. Uap → menggerakkan Turbin → listrik",
    "   4. Exhaust steam → proses (evaporator, vacuum pan, heater)",
    "",
    "▪ Data Kunci:",
    "   – 100 ton tebu → ~30 ton bagasse",
    "   – Bagasse = 100% bahan bakar boiler",
    "   – Surplus listrik → dijual ke PLN (3–5% revenue)",
    "",
    "▪ Utilitas Lainnya:",
    "   – Water treatment plant (air proses & imbibisi)",
    "   – Workshop & bengkel (maintenance)",
    "   – Laboratorium QC (analisis pol, brix, ICUMSA)",
]
add_body_text(slide12, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.5), bullets12, Pt(13))

# Energy cycle diagram (circular concept - simplified as vertical)
energy_items = ["Bagasse\n(Ampas Tebu)", "Boiler\n(Pembakaran)", "Uap Tekanan\nTinggi",
                "Turbin\n(Listrik)", "Exhaust Steam\n→ Proses"]
ey2 = Inches(1.8)
for i, lbl in enumerate(energy_items):
    c = DARK_GREEN if i in [0,2,4] else GOLD
    add_shape_box(slide12, Inches(9.2), ey2 + i * Inches(1.05), Inches(3.3), Inches(0.7),
                  c, lbl, Pt(11), WHITE, True)
    if i < len(energy_items) - 1:
        add_down_arrow(slide12, Inches(10.7), ey2 + i * Inches(1.05) + Inches(0.72))

add_key_takeaway(slide12, Inches(0.8), Inches(6.6), Inches(8),
                 "Pabrik gula = pembangkit listrik sendiri. Bagasse menggantikan bahan bakar fosil.")


# ============================================================
# SLIDE 13: BUSINESS/OPERATIONAL FLOW
# ============================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide13, WHITE)
add_title_bar(slide13, "ALUR BISNIS & OPERASIONAL",
              "Koordinasi, Logistik, Penjadwalan & Quality Control")

bullets13_left = [
    "▪ 6 Divisi Operasional:",
    "   Tanaman | Pabrikasi | Instalasi | Keuangan | SDM | Pemasaran",
    "",
    "▪ Koordinasi Operasional:",
    "   – Operasi 24/7 saat musim giling (shift 3×8 jam)",
    "   – Meeting harian: target giling vs aktual",
    "   – Komunikasi real-time antar stasiun",
    "",
    "▪ Logistik & Distribusi:",
    "   – Inbound: penjadwalan angkutan tebu dari kebun",
    "   – Outbound: pengiriman gula ke distributor/industri",
    "   – FIFO untuk inventori gudang",
    "",
    "▪ Quality Control:",
    "   – Lab menganalisis di setiap stasiun",
    "   – Parameter: pol, brix, ICUMSA, kadar air, warna",
    "   – Target losses keseluruhan: <2,5%",
]
add_body_text(slide13, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.5), bullets13_left, Pt(13))

# Monitoring layers
monitor_data = [
    ("HARIAN", "Kapasitas giling, rendemen,\nantrian tebu, breakdown", DARK_GREEN),
    ("MINGGUAN", "Target vs actual, K3,\nanalisis losses", MID_GREEN),
    ("BULANAN", "Cash flow, EBITDA,\nefisiensi energi", GOLD),
]
my2 = Inches(2.2)
for i, (period, desc, col) in enumerate(monitor_data):
    add_shape_box(slide13, Inches(9), my2 + i * Inches(1.6), Inches(3.8), Inches(1.2),
                  col, f"{period}\n{desc}", Pt(11), WHITE, True)

add_key_takeaway(slide13, Inches(0.8), Inches(6.6), Inches(8),
                 "Manajemen memonitor 3 layer: harian (operasional), mingguan (kualitas), bulanan (finansial)")

# ============================================================
# SLIDE 14: COMMON CHALLENGES
# ============================================================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide14, WHITE)
add_title_bar(slide14, "TANTANGAN OPERASIONAL PABRIK GULA",
              "Hambatan Industri & Bottleneck Produksi")

# Two columns: Operational Bottlenecks & Industry Challenges
col14_left = [
    "⚠️ BOTTLENECK OPERASIONAL:",
    "",
    "▪ Kualitas tebu rendah (pol <8%)",
    "▪ Breakdown mesin gilingan/boiler",
    "▪ Antrian tebu terlalu lama (pol turun)",
    "▪ Cuaca (hujan ganggu tebang/angkut)",
    "▪ SDM operator terbatas",
    "▪ Kapasitas terbatas vs pasokan",
    "",
    "Mitigasi: preventive maintenance,",
    "QC ketat di emplasement, training SDM",
]
col14_right = [
    "⚠️ TANTANGAN INDUSTRI:",
    "",
    "▪ Produktivitas rendah (RI 6-8% vs Thailand 10-12%)",
    "▪ Impor tinggi (~50% kebutuhan nasional)",
    "▪ Teknologi tua (mesin >30 tahun)",
    "▪ Margin tipis (HPP tinggi vs harga regulasi)",
    "▪ Perubahan iklim (kemarau/banjir)",
    "▪ Regenerasi petani tebu menurun",
    "",
    "Akar masalah: PRODUKTIVITAS",
    "→ itulah sebab impor tetap tinggi",
]
add_body_text(slide14, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2), col14_left, Pt(13))
add_body_text(slide14, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), col14_right, Pt(13))

# Opportunity boxes at bottom
opp_labels = [
    "Digitalisasi\n(IoT/SCADA/AI)\n+10-15% efisiensi",
    "Modernisasi\nMesin\n+1-2% rendemen",
    "Co-generation\nListrik\n+5-10% revenue",
    "Diversifikasi\n(Bioetanol, Refined)\n+15-25% revenue",
]
ox = Inches(0.6)
for i, lbl in enumerate(opp_labels):
    add_shape_box(slide14, ox + i * Inches(3.2), Inches(6.2), Inches(3.0), Inches(1.1),
                  SUGARCANE, lbl, Pt(10), WHITE, True)


# ============================================================
# SLIDE 15: KEY TAKEAWAYS / SUMMARY
# ============================================================
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide15, WHITE)
add_title_bar(slide15, "RANGKUMAN & KEY TAKEAWAYS",
              "5 Poin Penting yang Harus Diingat")

takeaways = [
    ("1", "Bisnis Terintegrasi", "Pabrik gula = hulu (kebun) + proses (pabrik) + hilir (distribusi). Satu rantai nilai yang saling tergantung.", DARK_GREEN),
    ("2", "Multi-Revenue Stream", "Bukan hanya gula — molasses, bagasse, blotong, dan listrik menyumbang 15–25% revenue tambahan.", MID_GREEN),
    ("3", "KPI Raja = Rendemen", "Target 8–10%. Kenaikan 1% rendemen = peningkatan profit signifikan. Semua stasiun berkontribusi.", SUGARCANE),
    ("4", "Tantangan Nyata", "Produktivitas rendah, teknologi tua, margin tipis. Indonesia masih impor 50% gula.", GOLD),
    ("5", "Peluang Besar", "Digitalisasi, modernisasi, co-gen, diversifikasi → potensi +20–30% EBITDA.", BROWN),
]

ty = Inches(1.6)
for i, (num, title, desc, color) in enumerate(takeaways):
    # Number circle
    add_shape_box(slide15, Inches(0.6), ty + i * Inches(1.1), Inches(0.6), Inches(0.6),
                  color, num, Pt(18), WHITE, True)
    # Title + description
    txb = slide15.shapes.add_textbox(Inches(1.4), ty + i * Inches(1.1), Inches(11), Inches(0.9))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_TITLE
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_TEXT
    p2.font.name = FONT_BODY

# Bottom quote
quote_box = slide15.shapes.add_textbox(Inches(1), Inches(7.0), Inches(11), Inches(0.4))
qtf = quote_box.text_frame
qp = qtf.paragraphs[0]
qp.text = '"100 ton tebu → 8–10 ton gula + 30 ton bagasse + 4–5 ton molasses = bisnis yang utuh"'
qp.font.size = Pt(13)
qp.font.italic = True
qp.font.color.rgb = GRAY
qp.font.name = FONT_BODY
qp.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 16: CLOSING SLIDE
# ============================================================
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide16, DARK_GREEN)

# Thank you text
txBox = slide16.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "TERIMA KASIH"
p.font.size = Pt(44)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = FONT_TITLE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "atas perhatian Bapak/Ibu"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(0xA5, 0xD6, 0xA7)
p2.font.name = FONT_BODY
p2.alignment = PP_ALIGN.CENTER

# Divider line
div = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.8), Inches(4), Inches(0.05))
div.fill.solid()
div.fill.fore_color.rgb = GOLD
div.line.fill.background()
div.shadow.inherit = False

# Q&A text
txBox2 = slide16.shapes.add_textbox(Inches(2), Inches(5.0), Inches(9), Inches(1.5))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "Sesi Tanya Jawab"
p3.font.size = Pt(22)
p3.font.color.rgb = GOLD
p3.font.bold = True
p3.font.name = FONT_TITLE
p3.alignment = PP_ALIGN.CENTER

p4 = tf2.add_paragraph()
p4.text = "Silakan ajukan pertanyaan atau berikan masukan"
p4.font.size = Pt(14)
p4.font.color.rgb = RGBColor(0xC8, 0xE6, 0xC9)
p4.font.name = FONT_BODY
p4.alignment = PP_ALIGN.CENTER

# ============================================================
# SAVE
# ============================================================
output_path = '/projects/sandbox/buat-tugas/Overview_Proses_Alur_Bisnis_Pabrik_Gula.pptx'
prs.save(output_path)
print(f"✅ Presentasi berhasil dibuat: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
